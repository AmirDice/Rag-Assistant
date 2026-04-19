"""WP13 — Format conversion: PDF / DOCX / PPTX / XLSX → Markdown.

Each converter returns a tuple of (markdown_text, list[extracted_image_paths]).
"""

from __future__ import annotations

import io
import logging
import re
import subprocess
import tempfile
from pathlib import Path
from typing import Protocol

from PIL import Image as PILImage

logger = logging.getLogger(__name__)

_PICTURE_MARKER_RE = re.compile(
    r"\*?\*?\s*==>\s*picture\s*\[\d+\s*x\s*\d+\]\s*intentionally omitted\s*<==\s*\*?\*?",
    re.IGNORECASE,
)

_EMPTY_BRACKET_RE = re.compile(r"\(\s*\)")

_ICONO_EMPTY_RE = re.compile(
    r"(icono|botón|botó|ícono|icon)\s*\(\s*\)",
    re.IGNORECASE,
)

_ICONO_BEFORE_BRACKET_RE = re.compile(
    r"(icono|ícono|icon)\s+\(\s*\)",
    re.IGNORECASE,
)

_REPEATED_CHARS_RE = re.compile(r"([a-záéíóúñ])\1{4,}", re.IGNORECASE)
_GARBLED_DOTS_RE = re.compile(r"\.{5,}")
_TOC_NOISE_RE = re.compile(
    r"[.oOcCnNrR ]{10,}\d+\s*$", re.MULTILINE
)
_REPEATED_WORDS_RE = re.compile(r"\b(\w{2,})\s+(?:\1\s+){2,}", re.IGNORECASE)
_MULTI_SPACE_RE = re.compile(r"[ \t]{3,}")
_HEADER_GARBLE_RE = re.compile(
    r"^(?:Far\s*v?M?At?T?I?i?c|FazmaA?r?I?i?c|LARM[AI]ATIC|LARvMIATIC)\s*$",
    re.MULTILINE | re.IGNORECASE,
)


def _clean_ocr_text(text: str) -> str:
    """Remove common OCR artifacts from pytesseract output."""
    text = _HEADER_GARBLE_RE.sub("", text)
    text = _TOC_NOISE_RE.sub("", text)
    text = _REPEATED_CHARS_RE.sub(r"\1\1", text)
    text = _GARBLED_DOTS_RE.sub("...", text)
    text = _REPEATED_WORDS_RE.sub(r"\1 ", text)
    text = _MULTI_SPACE_RE.sub(" ", text)
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            cleaned.append("")
            continue
        alpha = sum(1 for c in stripped if c.isalpha())
        if len(stripped) > 5 and alpha / len(stripped) < 0.15:
            continue
        cleaned.append(line)
    text = "\n".join(cleaned)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text


def _clean_pymupdf_markers(md: str) -> str:
    """Remove pymupdf4llm 'intentionally omitted' placeholders and fix empty brackets."""
    md = _PICTURE_MARKER_RE.sub("", md)

    md = _ICONO_EMPTY_RE.sub(r"\1 (icono)", md)

    def _replace_empty_bracket(m: re.Match) -> str:
        start = max(0, m.start() - 40)
        context_before = md[start:m.start()].lower()
        if any(kw in context_before for kw in ("clic", "pulsar", "presionar", "click", "pinchar")):
            return "(botón)"
        return "(icono)"

    md = _EMPTY_BRACKET_RE.sub(_replace_empty_bracket, md)
    return md


class Converter(Protocol):
    def convert(self, source: Path, output_dir: Path) -> tuple[str, list[Path]]: ...


# ── PDF (native text) via pymupdf4llm ────────────────────────

class PDFConverter:
    def convert(self, source: Path, output_dir: Path) -> tuple[str, list[Path]]:
        import pymupdf4llm
        import pymupdf

        doc = pymupdf.open(str(source))
        images: list[Path] = []
        media_dir = output_dir / "media"

        sparse_pages = []
        for i, page in enumerate(doc):
            if len(page.get_text().strip()) < 100:
                sparse_pages.append(i)

        if len(sparse_pages) > len(doc) * 0.5:
            logger.info(
                "OCR needed for %s — %d/%d pages sparse, using pytesseract",
                source.name, len(sparse_pages), len(doc),
            )
            md_text = self._ocr_convert(source)
        elif sparse_pages:
            logger.info(
                "Hybrid extraction for %s — OCR on %d sparse pages, native on %d",
                source.name, len(sparse_pages), len(doc) - len(sparse_pages),
            )
            md_text = self._hybrid_convert(source, set(sparse_pages))
        else:
            md_text = pymupdf4llm.to_markdown(str(source))

        md_text = _clean_pymupdf_markers(md_text)

        from api.pipeline.logo_filter import is_logo_pil

        media_dir.mkdir(parents=True, exist_ok=True)
        for page_num, page in enumerate(doc):
            if page_num == 0:
                continue

            page_width = page.rect.width
            for img_idx, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                try:
                    pix = pymupdf.Pixmap(doc, xref)
                    if pix.width < 15 or pix.height < 10:
                        continue

                    if pix.width == 20 and pix.height == 29:
                        continue

                    if pix.n - pix.alpha > 3:
                        pix = pymupdf.Pixmap(pymupdf.csRGB, pix)

                    png_bytes = pix.tobytes("png")
                    pil_img = PILImage.open(io.BytesIO(png_bytes))
                    if is_logo_pil(pil_img):
                        logger.debug("Skipping logo image xref=%d on page %d", xref, page_num)
                        pil_img.close()
                        continue
                    pil_img.close()

                    is_small = pix.width < 60 and pix.height < 60
                    if is_small:
                        rects = page.get_image_rects(xref)
                        if rects:
                            x0 = rects[0].x0
                            if x0 < page_width * 0.15:
                                continue

                    img_name = f"image_p{page_num + 1:03d}_{img_idx:03d}.png"
                    img_path = media_dir / img_name
                    pix.save(str(img_path))
                    images.append(img_path)
                except Exception as e:
                    logger.warning("Failed to extract image xref=%d: %s", xref, e)
        doc.close()
        return md_text, images

    def _ocr_convert(self, source: Path) -> str:
        try:
            import pytesseract
            from PIL import Image
            import pymupdf

            doc = pymupdf.open(str(source))
            pages = []
            for page_num, page in enumerate(doc):
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img, lang="spa")
                text = _clean_ocr_text(text)
                if text.strip():
                    pages.append(f"## Página {page_num + 1}\n\n{text}")
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            logger.error("pytesseract not installed — OCR skipped")
            return ""

    def _hybrid_convert(self, source: Path, ocr_pages: set[int]) -> str:
        """Use pymupdf4llm for native-text pages, pytesseract for sparse ones."""
        import pymupdf4llm
        import pymupdf

        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            logger.warning("pytesseract not available, using pymupdf4llm for all pages")
            return pymupdf4llm.to_markdown(str(source))

        native_page_list = sorted(
            i for i in range(pymupdf.open(str(source)).page_count)
            if i not in ocr_pages
        )
        native_md = pymupdf4llm.to_markdown(str(source), pages=native_page_list) if native_page_list else ""

        ocr_parts: list[str] = []
        if ocr_pages:
            doc = pymupdf.open(str(source))
            for page_num in sorted(ocr_pages):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img, lang="spa")
                text = _clean_ocr_text(text)
                if text.strip():
                    ocr_parts.append(f"## Página {page_num + 1}\n\n{text}")
            doc.close()

        if ocr_parts:
            return native_md + "\n\n" + "\n\n".join(ocr_parts)
        return native_md


# ── DOCX via python-docx + pandoc ────────────────────────────

class DOCXConverter:
    def convert(self, source: Path, output_dir: Path) -> tuple[str, list[Path]]:
        images = self._extract_images(source, output_dir)
        md_text = self._pandoc_convert(source)
        return md_text, images

    def _pandoc_convert(self, source: Path) -> str:
        try:
            result = subprocess.run(
                ["pandoc", str(source), "-t", "markdown", "--wrap=none"],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0:
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired) as e:
            logger.warning("pandoc failed: %s — falling back to python-docx", e)

        return self._fallback_convert(source)

    def _fallback_convert(self, source: Path) -> str:
        import docx
        doc = docx.Document(str(source))
        parts = []
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if "Heading 1" in style:
                parts.append(f"# {text}")
            elif "Heading 2" in style:
                parts.append(f"## {text}")
            elif "Heading 3" in style:
                parts.append(f"### {text}")
            else:
                parts.append(text)
        return "\n\n".join(parts)

    def _extract_images(self, source: Path, output_dir: Path) -> list[Path]:
        import docx
        from docx.opc.constants import RELATIONSHIP_TYPE as RT

        images: list[Path] = []
        media_dir = output_dir / "media"
        media_dir.mkdir(parents=True, exist_ok=True)

        doc = docx.Document(str(source))
        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    ext = Path(rel.target_ref).suffix or ".png"
                    img_path = media_dir / f"image_{i:03d}{ext}"
                    img_path.write_bytes(blob)
                    images.append(img_path)
                except Exception as e:
                    logger.warning("Failed to extract DOCX image: %s", e)
        return images


# ── PPTX: slides → PDF via LibreOffice → marker ─────────────

class PPTXConverter:
    def convert(self, source: Path, output_dir: Path) -> tuple[str, list[Path]]:
        # Extract speaker notes separately
        notes = self._extract_notes(source)

        # Convert to PDF via LibreOffice, then to Markdown
        md_text = self._via_libreoffice(source, output_dir)
        if notes:
            md_text += "\n\n## Notas del presentador\n\n" + notes

        return md_text, []  # images handled via PDF conversion

    def _via_libreoffice(self, source: Path, output_dir: Path) -> str:
        with tempfile.TemporaryDirectory() as tmpdir:
            try:
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pdf",
                     "--outdir", tmpdir, str(source)],
                    capture_output=True, timeout=120,
                )
                pdf_path = Path(tmpdir) / (source.stem + ".pdf")
                if pdf_path.exists():
                    pdf_conv = PDFConverter()
                    md, _ = pdf_conv.convert(pdf_path, output_dir)
                    return md
            except (FileNotFoundError, subprocess.TimeoutExpired) as e:
                logger.warning("LibreOffice conversion failed: %s", e)

        return self._fallback_text(source)

    def _extract_notes(self, source: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(source))
            notes = []
            for i, slide in enumerate(prs.slides, 1):
                if slide.has_notes_slide:
                    text = slide.notes_slide.notes_text_frame.text.strip()
                    if text:
                        notes.append(f"**Diapositiva {i}:** {text}")
            return "\n\n".join(notes)
        except Exception as e:
            logger.warning("Failed to extract PPTX notes: %s", e)
            return ""

    def _fallback_text(self, source: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(source))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"## Diapositiva {i}\n\n" + "\n\n".join(texts))
            return "\n\n".join(parts)
        except Exception:
            return ""


# ── XLSX via openpyxl → Markdown tables ──────────────────────

class XLSXConverter:
    def convert(self, source: Path, output_dir: Path) -> tuple[str, list[Path]]:
        import openpyxl
        wb = openpyxl.load_workbook(str(source), data_only=True)
        sections = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            section = f"## {sheet_name}\n\n"
            # Build markdown table
            header = rows[0]
            cols = [str(c) if c is not None else "" for c in header]
            section += "| " + " | ".join(cols) + " |\n"
            section += "| " + " | ".join("---" for _ in cols) + " |\n"
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]
                section += "| " + " | ".join(cells) + " |\n"
            sections.append(section)

        return "\n\n".join(sections), []


# ── Factory ──────────────────────────────────────────────────

_CONVERTERS: dict[str, Converter] = {
    "pdf": PDFConverter(),
    "docx": DOCXConverter(),
    "pptx": PPTXConverter(),
    "xlsx": XLSXConverter(),
}


def get_converter(fmt: str) -> Converter:
    fmt = fmt.lower().lstrip(".")
    if fmt not in _CONVERTERS:
        raise ValueError(f"Unsupported format: {fmt}. Supported: {list(_CONVERTERS)}")
    return _CONVERTERS[fmt]


def convert_document(source: Path, output_dir: Path) -> tuple[str, list[Path]]:
    """Convert any supported document to Markdown + extracted images."""
    fmt = source.suffix.lower().lstrip(".")
    converter = get_converter(fmt)
    output_dir.mkdir(parents=True, exist_ok=True)
    logger.info("Converting %s (format=%s)", source.name, fmt)
    return converter.convert(source, output_dir)
